#!/usr/bin/env python3
# 生成最终版 PPTX 并尝试自动导出 PDF
# 依赖：python-pptx
# pip install python-pptx

from pptx import Presentation
from pptx.util import Cm, Pt, Inches
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
import subprocess, os, sys

pptx_filename = "大学生职业生涯规划书.pptx"
pdf_filename = "大学生职业生涯规划书.pdf"

# 幻灯片内容（同上最终文本）
slides = [
    {"title":"大学生职业生涯规划书",
     "lines":["姓名：XXX","学号：2024352XXXX","班级：1班","指导教师：JIDAUN LOUDAFENG","提交日期：2026-06-19"],
     "notes":"封面。"},
    {"title":"职业生涯彩虹图",
     "lines":["覆盖成长、探索、建立、维持、退出五大阶段","六大角色：学生、工作者、休闲者、公民、持家者、子女","结合 ISTP 与嵌入式职业赛道的全生命周期规划"],
     "notes":"说明个人特质如何映射到彩虹图各阶段。"},
    {"title":"自我认知 — 职业兴趣",
     "lines":["霍兰德测评：RIE（现实型/研究型/企业型）","偏好：硬件实操、代码调试、设备联调；坚持跑步","兴趣与嵌入式软硬结合方向高度匹配"],
     "notes":"举例兴趣在项目/竞赛中的体现。"},
    {"title":"自我认知 — 职业性格",
     "lines":["MBTI：ISTP（手艺人型）","优势：理性务实、动手能力强、冷静沉稳、抗压性好","短板：公开汇报与长篇写作需加强"],
     "notes":"提出通过训练与任务分工提升短板的办法。"},
    {"title":"自我认知 — 职业价值观",
     "lines":["三大价值：技术成长、工作自主度、简洁技术团队氛围","嵌入式岗位与上述价值高度契合"],
     "notes":"面试或择岗时的优先级说明。"},
    {"title":"自我认知 — 职业能力（优势）",
     "lines":["逻辑统筹：校青协干事经历","沟通协作：活动组织与竞赛配合经验","实操自学：实验室与硬件调试经验","抗压持久力：���期跑步培养的耐力","目标明确：计划考级、参与竞赛、入实验室"],
     "notes":"每项配对应证明或行动举例。"},
    {"title":"自我认知 — 职业能力（弱势）",
     "lines":["底层理论需完善（OS/编译原理/微机原理）","大型项目架构经验不足","公开演讲与长文写作能力有待提升"],
     "notes":"补短策略：课程学习、导师辅导、反复演练。"},
    {"title":"决策平衡单",
     "lines":["两条路径对比：嵌入式 vs 后端","权重打分后结论：嵌入式得分更高 → 结论：嵌入式"],
     "notes":"说明关键权重与选择理由。"},
    {"title":"家庭 / 学校 / 社会环境分析",
     "lines":["家庭：经济支持、鼓励技术路线","学校：课程与实验室资源（STM32/ARM/Linux 等）","社会：物联网/车用电子等领域对嵌入式人才需求大"],
     "notes":"强调资源如何支持短中期目标实现。"},
    {"title":"目标职业：嵌入式开发工程师",
     "lines":["岗位职责：底层驱动、软硬联调、性能优化、量产支持","任职要求：C/C++、单片机/ARM/STM32、通信协议、Linux 基础"],
     "notes":"说明个人优势如何匹配岗位需求。"},
    {"title":"SWOT 分析",
     "lines":["优势：ISTP、青协、抗压、自学、目标明确","劣势：硬件电路、项目架构、汇报能力","机会：行业需求、学校资源、在线学习","威胁：毕业生增多、研究生门槛、项目经验要求"],
     "notes":"对应策略：补短、利用资源、竞赛/实验室。"},
    {"title":"职业目标分解（短/中/长）",
     "lines":["短期：四六级、入实验室、竞赛、作品集","中期：研究生/企业实习、独立项目","长期：高级工程师、方案设计、技术带队"],
     "notes":"每阶段关键指标与时间节点。"},
    {"title":"计划与评估调整 / 结束语",
     "lines":["评估频率：学期复盘 / 年度复盘","调整原则：实习优先、备选岗位、动态复盘","名言：道阻且长，行则将至；行而不辍，未来可期。"],
     "notes":"鼓励与下一步行动建议。"}
]

prs = Presentation()
# 设置幻灯片大小（默认），并计算页边距
slide_w = prs.slide_width
slide_h = prs.slide_height
margin = Cm(2.5)
content_w = slide_w - margin*2
content_h = slide_h - margin*2

# 字体对应（尽量兼容中/英文系统）
TITLE_FONT = "SimHei"   # 黑体（SimHei）
BODY_FONT = "SimSun"    # 宋体（SimSun）

# 字号近似值（pt）
SIZE_TITLE_YIJI = Pt(22)   # 二号近似 22pt
SIZE_TITLE_ERJI = Pt(16)   # 小三约 16pt
SIZE_BODY = Pt(12)         # 小四约 12pt
SIZE_NOTE = Pt(10.5)       # 五号约 10.5pt

for idx, s in enumerate(slides):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    # title textbox
    title_h = Cm(2.5)
    tb = slide.shapes.add_textbox(margin, margin, content_w, title_h)
    tf = tb.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = s["title"]
    p.font.name = TITLE_FONT
    p.font.size = SIZE_TITLE_YIJI
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    # body textbox
    body_top = margin + title_h + Cm(0.5)
    body_box = slide.shapes.add_textbox(margin, body_top, content_w, content_h - title_h - Cm(0.5))
    body_tf = body_box.text_frame
    body_tf.clear()
    body_tf.word_wrap = True
    # 添加每行为一个段落，设置段后、首行缩进、行距
    for i, line in enumerate(s["lines"]):
        p = body_tf.add_paragraph()
        p.text = line
        p.level = 0
        p.font.name = BODY_FONT
        p.font.size = SIZE_BODY
        p.alignment = PP_ALIGN.LEFT
        # 首行缩进（约2字符）
        try:
            p.paragraph_format.first_line_indent = Pt(24)  # 24pt ~ 两个中文字符
            p.paragraph_format.space_after = Pt(6)         # 段后间距约 6pt（接近 1 行的视觉间距）
            p.line_spacing = 1.5
        except Exception:
            pass
    # 如果需要图表注释，可以在此处添加（示例不包含图）
    # 备注（speaker notes）
    notes_slide = slide.notes_slide
    notes_tf = notes_slide.notes_text_frame
    notes_tf.text = s.get("notes", "")

# 删除默认最初空幻灯片（如果存在）
if prs.slides:
    # first slide in layout is actually blank added by Presentation() - but we added slides after it, so keep as-is.

    pass

# 保存 pptx
prs.save(pptx_filename)
print("Saved PPTX:", pptx_filename)

# 尝试使用 libreoffice 导出为 PDF
def convert_to_pdf(pptx, out_pdf):
    try:
        subprocess.run(["libreoffice", "--version"], check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
    except Exception as e:
        print("LibreOffice 未检测到，无法自动导出 PDF。请在本地使用 LibreOffice 或 PowerPoint 将 PPTX 导出为 PDF。")
        return False
    try:
        subprocess.run(["libreoffice", "--headless", "--convert-to", "pdf", pptx, "--outdir", "."], check=True)
        base_pdf = os.path.splitext(pptx)[0] + ".pdf"
        if os.path.exists(base_pdf):
            if base_pdf != out_pdf:
                os.replace(base_pdf, out_pdf)
            print("生成 PDF:", out_pdf)
            return True
    except Exception as e:
        print("LibreOffice 导出 PDF 失败：", e)
    return False

ok = convert_to_pdf(pptx_filename, pdf_filename)
if not ok:
    print("请手动将 %s 转换为 PDF（使用 PowerPoint/LibreOffice）。" % pptx_filename)
else:
    print("生成完成：", pptx_filename, pdf_filename)
